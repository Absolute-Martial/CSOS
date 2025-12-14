"""
Personal Engineering OS - File Handler
Extract text from PDF, DOCX, PPTX for AI reading
"""

import os
import re
from pathlib import Path
from typing import Optional, Tuple

import aiofiles


# Supported file extensions
SUPPORTED_EXTENSIONS = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain",
    ".md": "text/markdown"
}

UPLOAD_DIR = os.getenv("UPLOAD_DIR", "./uploads")


def validate_filename(filename: str) -> Tuple[bool, str]:
    """Validate filename follows naming convention."""
    # Pattern: snake_case with extension
    pattern = r"^[a-z][a-z0-9_]*\.[a-z]+$"
    
    if not re.match(pattern, filename.lower()):
        return False, "Filename must be snake_case (e.g., lecture_notes.pdf)"
    
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return False, f"Unsupported file type. Allowed: {', '.join(SUPPORTED_EXTENSIONS.keys())}"
    
    return True, "Valid"


def get_folder_path(subject_code: str, chapter_number: int, file_type: str) -> str:
    """Get the folder path for a file."""
    chapter_folder = f"chapter{chapter_number:02d}"
    return os.path.join(UPLOAD_DIR, subject_code.upper(), chapter_folder, file_type)


async def save_uploaded_file(
    content: bytes,
    subject_code: str,
    chapter_number: int,
    file_type: str,  # 'slides', 'assignments', 'notes'
    filename: str
) -> Tuple[str, int]:
    """Save file to appropriate folder, return path and size."""
    
    folder_path = get_folder_path(subject_code, chapter_number, file_type)
    os.makedirs(folder_path, exist_ok=True)
    
    file_path = os.path.join(folder_path, filename)
    
    async with aiofiles.open(file_path, 'wb') as f:
        await f.write(content)
    
    return file_path, len(content)


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file."""
    try:
        from pypdf import PdfReader
        
        reader = PdfReader(file_path)
        text_parts = []
        
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"Error extracting PDF text: {str(e)}"


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file."""
    try:
        from docx import Document
        
        doc = Document(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"Error extracting DOCX text: {str(e)}"


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file."""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"Error extracting PPTX text: {str(e)}"


async def extract_text_from_txt(file_path: str) -> str:
    """Read text from TXT or MD file."""
    try:
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            return await f.read()
    except Exception as e:
        return f"Error reading file: {str(e)}"


async def read_file_content(file_path: str) -> str:
    """Extract text from any supported file type."""
    ext = Path(file_path).suffix.lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext in [".txt", ".md"]:
        return await extract_text_from_txt(file_path)
    else:
        return f"Unsupported file type: {ext}"


def get_file_info(file_path: str) -> dict:
    """Get file metadata."""
    path = Path(file_path)
    
    if not path.exists():
        return {"error": "File not found"}
    
    stat = path.stat()
    
    return {
        "filename": path.name,
        "extension": path.suffix,
        "size_bytes": stat.st_size,
        "size_human": format_size(stat.st_size),
        "mimetype": SUPPORTED_EXTENSIONS.get(path.suffix.lower(), "application/octet-stream"),
        "modified": stat.st_mtime
    }


def format_size(size_bytes: int) -> str:
    """Format bytes to human readable."""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.1f} TB"


def list_chapter_files(subject_code: str, chapter_number: int) -> dict:
    """List all files in a chapter folder."""
    base_path = os.path.join(UPLOAD_DIR, subject_code.upper(), f"chapter{chapter_number:02d}")
    
    if not os.path.exists(base_path):
        return {"error": "Chapter folder not found", "path": base_path}
    
    files = {
        "slides": [],
        "assignments": [],
        "notes": []
    }
    
    for file_type in files.keys():
        folder = os.path.join(base_path, file_type)
        if os.path.exists(folder):
            for filename in os.listdir(folder):
                file_path = os.path.join(folder, filename)
                if os.path.isfile(file_path):
                    files[file_type].append(get_file_info(file_path))
    
    return {
        "subject": subject_code,
        "chapter": chapter_number,
        "path": base_path,
        "files": files
    }
